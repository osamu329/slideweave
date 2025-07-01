#!/usr/bin/env python3
# /// script
# requires-python = ">=3.8"
# dependencies = [
#   "python-pptx>=0.6.21",
# ]
# ///
"""
PowerPointの解像度とDPI設定を調査するスクリプト
"""

from pptx import Presentation
from pptx.util import Inches

def create_test_presentation():
    """テスト用のプレゼンテーションを作成"""
    
    # 新しいプレゼンテーションを作成
    prs = Presentation()
    
    # デフォルトスライドサイズを確認
    print("=== デフォルトスライドサイズ ===")
    print(f"幅: {prs.slide_width / 914400:.3f} インチ")
    print(f"高さ: {prs.slide_height / 914400:.3f} インチ")
    print(f"幅(EMU): {prs.slide_width}")
    print(f"高さ(EMU): {prs.slide_height}")
    
    # ピクセル換算（96DPI）
    width_px_96 = (prs.slide_width / 914400) * 96
    height_px_96 = (prs.slide_height / 914400) * 96
    print(f"96DPI換算: {width_px_96:.0f} × {height_px_96:.0f} ピクセル")
    
    # ピクセル換算（72DPI）
    width_px_72 = (prs.slide_width / 914400) * 72
    height_px_72 = (prs.slide_height / 914400) * 72
    print(f"72DPI換算: {width_px_72:.0f} × {height_px_72:.0f} ピクセル")
    
    # Widescreenサイズに設定（16:9）
    widescreen_width = Inches(13.333)
    widescreen_height = Inches(7.5)
    
    prs.slide_width = widescreen_width
    prs.slide_height = widescreen_height
    
    print("\n=== Widescreenサイズ設定後 ===")
    print(f"幅: {prs.slide_width / 914400:.3f} インチ")
    print(f"高さ: {prs.slide_height / 914400:.3f} インチ")
    print(f"幅(EMU): {prs.slide_width}")
    print(f"高さ(EMU): {prs.slide_height}")
    
    # ピクセル換算（96DPI）
    width_px_96 = (prs.slide_width / 914400) * 96
    height_px_96 = (prs.slide_height / 914400) * 96
    print(f"96DPI換算: {width_px_96:.0f} × {height_px_96:.0f} ピクセル")
    
    # ピクセル換算（72DPI）
    width_px_72 = (prs.slide_width / 914400) * 72
    height_px_72 = (prs.slide_height / 914400) * 72
    print(f"72DPI換算: {width_px_72:.0f} × {height_px_72:.0f} ピクセル")
    
    # スライドを追加
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # テスト用テキストボックスを追加
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    text_frame = textbox.text_frame
    text_frame.text = "Resolution Test Slide"
    
    # 保存
    prs.save('examples/output/resolution_test.pptx')
    print(f"\nテストファイルを保存しました: examples/output/resolution_test.pptx")
    
    return prs

def analyze_dpi_relationships():
    """DPIの関係性を分析"""
    print("\n=== DPI関係性分析 ===")
    
    # 13.333" × 7.5" Widescreen の場合
    width_inches = 13.333
    height_inches = 7.5
    
    print(f"Widescreen サイズ: {width_inches}\" × {height_inches}\"")
    print(f"アスペクト比: {width_inches/height_inches:.3f}:1 (16:9)")
    
    # 各DPIでのピクセル値
    dpis = [72, 96, 144, 192, 288]
    
    for dpi in dpis:
        width_px = width_inches * dpi
        height_px = height_inches * dpi
        print(f"{dpi:3d} DPI: {width_px:4.0f} × {height_px:3.0f} ピクセル")
        
        # 標準解像度との比較
        if width_px == 960 and height_px == 540:
            print(f"         → これが 960×540 の基準 (72 DPI)")
        elif width_px == 1280 and height_px == 720:
            print(f"         → これが 1280×720 の基準 (96 DPI, HD 720p)")
        elif width_px == 1920 and height_px == 1080:
            print(f"         → これが 1920×1080 の基準 (144 DPI, Full HD)")

def explain_powerpoint_resolution():
    """PowerPointの解像度システムを説明"""
    print("\n=== PowerPoint解像度システムの説明 ===")
    
    print("1. PowerPointの内部単位系:")
    print("   - 基本単位: EMU (English Metric Units)")
    print("   - 1インチ = 914,400 EMU")
    print("   - 1ポイント = 12,700 EMU (1/72インチ)")
    
    print("\n2. Widescreenスライド (16:9):")
    print("   - 設定サイズ: 13.333\" × 7.5\"")
    print("   - これは意図的に16:9の比率を保つための値")
    print("   - 13.333 / 7.5 = 1.7777... ≈ 16/9")
    
    print("\n3. 表示解像度とDPIの関係:")
    print("   - 72 DPI: 960×540 ピクセル (Typography標準)")
    print("   - 96 DPI: 1280×720 ピクセル (Windows標準)")
    print("   - 編集時の表示は通常96DPIベース")
    print("   - エクスポート時は設定可能")
    
    print("\n4. なぜ960×540が表示されるのか:")
    print("   - PowerPointの内部計算で72DPIが基準として使用される")
    print("   - 13.333 × 72 = 960 ピクセル")
    print("   - 7.5 × 72 = 540 ピクセル")
    print("   - これがUI上で「実際の解像度」として表示される")

if __name__ == "__main__":
    create_test_presentation()
    analyze_dpi_relationships()
    explain_powerpoint_resolution()