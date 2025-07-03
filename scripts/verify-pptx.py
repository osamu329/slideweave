#!/usr/bin/env python3
# /// script
# requires-python = ">=3.8"
# dependencies = [
#   "python-pptx>=0.6.21",
# ]
# ///
"""
PPTXファイルの内容を検証するスクリプト
python-pptxを使用してPPTXの実際の内容を解析し、レポートを生成

使用方法:
  uv run scripts/verify-pptx.py <pptx-file>
"""

import sys
import json
import re
import base64
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_svg_info(image_blob):
    """SVG画像データから全ての情報を抽出"""
    try:
        svg_content = None
        # 画像データをテキストとしてデコード（SVGの場合）
        if image_blob.startswith(b'<svg') or b'<svg' in image_blob[:100]:
            svg_content = image_blob.decode('utf-8')
        # Data URIの場合（base64エンコード済みSVG）
        elif image_blob.startswith(b'data:image/svg+xml;base64,'):
            base64_data = image_blob[26:]  # プレフィックスを除去
            svg_content = base64.b64decode(base64_data).decode('utf-8')
        
        if svg_content:
            info = {}
            # fill属性
            fill_match = re.search(r'fill="([^"]+)"', svg_content)
            if fill_match:
                info['fill'] = fill_match.group(1)
            
            # width/height属性
            width_match = re.search(r'width="([^"]+)"', svg_content)
            height_match = re.search(r'height="([^"]+)"', svg_content)
            if width_match:
                info['width'] = width_match.group(1)
            if height_match:
                info['height'] = height_match.group(1)
            
            # stroke関連
            stroke_match = re.search(r'stroke="([^"]+)"', svg_content)
            if stroke_match:
                info['stroke'] = stroke_match.group(1)
            
            stroke_width_match = re.search(r'stroke-width="([^"]+)"', svg_content)
            if stroke_width_match:
                info['stroke_width'] = stroke_width_match.group(1)
                
            # opacity
            opacity_match = re.search(r'opacity="([^"]+)"', svg_content)
            if opacity_match:
                info['opacity'] = opacity_match.group(1)
                
            # rect要素の確認
            rect_matches = re.findall(r'<rect[^>]*>', svg_content)
            if rect_matches:
                info['rect_count'] = len(rect_matches)
                info['rects'] = rect_matches
            
            # SVG全体のサイズが0でないか確認
            if svg_content.strip() == '' or len(svg_content) < 50:
                info['error'] = 'SVG content too short'
            
            return info
    except Exception as e:
        return {'error': str(e)}
    return None

def analyze_shape(shape, indent=0):
    """シェイプの情報を解析"""
    prefix = "  " * indent
    info = {
        "type": shape.shape_type,
        "name": shape.name,
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }
    
    # 位置とサイズ（EMU to インチ変換）
    print(f"{prefix}Shape: {shape.name}")
    print(f"{prefix}  Position: ({shape.left/914400:.2f}\", {shape.top/914400:.2f}\")")
    print(f"{prefix}  Size: {shape.width/914400:.2f}\" x {shape.height/914400:.2f}\"")
    
    # 背景色の確認
    if hasattr(shape, 'fill'):
        fill = shape.fill
        if fill.type == 1:  # Solid fill
            if fill.fore_color.rgb:
                print(f"{prefix}  Fill Color: #{fill.fore_color.rgb}")
            info["fill_color"] = str(fill.fore_color.rgb) if fill.fore_color.rgb else None
    
    # SVG画像の場合、全情報を抽出
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            # 画像データからSVGの全情報を抽出
            image = shape.image
            if image:
                # SVGデータから全情報を抽出（SlideWeave生成の場合）
                svg_info = extract_svg_info(image.blob)
                if svg_info:
                    print(f"{prefix}  SVG Info:")
                    for key, value in svg_info.items():
                        if key == 'rects':
                            print(f"{prefix}    {key}: {len(value)} rect elements")
                            for i, rect in enumerate(value):
                                print(f"{prefix}      rect {i}: {rect}")
                        else:
                            print(f"{prefix}    {key}: {value}")
                    info["svg_info"] = svg_info
        except Exception as e:
            print(f"{prefix}  SVG Error: {e}")
            pass  # 画像の場合はスキップ
    
    # テキストの確認
    if shape.has_text_frame:
        text_frame = shape.text_frame
        
        # テキストフレームの垂直位置設定
        if hasattr(text_frame, 'vertical_anchor'):
            anchor_map = {0: "top", 1: "middle", 2: "bottom", 3: "mixed"}
            anchor = anchor_map.get(text_frame.vertical_anchor, "unknown")
            print(f"{prefix}  Vertical Anchor: {anchor}")
            info["vertical_anchor"] = anchor
        
        # テキストフレームのマージン設定（paddingに相当）
        if hasattr(text_frame, 'margin_left'):
            margin_left = text_frame.margin_left / 914400 if text_frame.margin_left else 0
            margin_top = text_frame.margin_top / 914400 if text_frame.margin_top else 0
            margin_right = text_frame.margin_right / 914400 if text_frame.margin_right else 0
            margin_bottom = text_frame.margin_bottom / 914400 if text_frame.margin_bottom else 0
            print(f"{prefix}  Text Margins (inches): L:{margin_left:.3f} T:{margin_top:.3f} R:{margin_right:.3f} B:{margin_bottom:.3f}")
            info["text_margin_left"] = margin_left
            info["text_margin_top"] = margin_top
            info["text_margin_right"] = margin_right
            info["text_margin_bottom"] = margin_bottom
        
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    print(f"{prefix}  Text: \"{run.text}\"")
                    try:
                        if run.font.color.rgb:
                            print(f"{prefix}    Font Color: #{run.font.color.rgb}")
                    except AttributeError:
                        print(f"{prefix}    Font Color: default/none")
                    font_size = run.font.size/12700 if run.font.size else "default"
                    print(f"{prefix}    Font Size: {font_size}pt" if isinstance(font_size, float) else f"{prefix}    Font Size: {font_size}")
                    print(f"{prefix}    Bold: {run.font.bold}")
                    print(f"{prefix}    Italic: {run.font.italic}")
                    if run.font.name:
                        print(f"{prefix}    Font Family: {run.font.name}")
                    info["text"] = run.text
                    try:
                        info["font_color"] = str(run.font.color.rgb) if run.font.color.rgb else None
                    except AttributeError:
                        info["font_color"] = None
                    info["font_size"] = run.font.size/12700 if run.font.size else None
                    info["font_bold"] = run.font.bold
                    info["font_italic"] = run.font.italic
                    info["font_family"] = run.font.name
    
    # グループシェイプの場合、子要素を再帰的に解析
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        print(f"{prefix}  Group with {len(shape.shapes)} shapes:")
        for child_shape in shape.shapes:
            analyze_shape(child_shape, indent + 1)
    
    return info

def verify_pptx(pptx_path):
    """PPTXファイルを検証"""
    print(f"=== PPTX Verification: {pptx_path} ===\n")
    
    prs = Presentation(pptx_path)
    
    # プレゼンテーション情報
    print(f"Slide Size: {prs.slide_width/914400:.1f}\" x {prs.slide_height/914400:.1f}\"")
    print(f"Number of Slides: {len(prs.slides)}\n")
    
    all_slides_info = []
    
    # 各スライドを解析
    for idx, slide in enumerate(prs.slides):
        print(f"=== Slide {idx + 1} ===")
        slide_info = {
            "slide_number": idx + 1,
            "shapes": []
        }
        
        # 背景色の確認
        if slide.background:
            if hasattr(slide.background, 'fill'):
                fill = slide.background.fill
                if fill.type == 1 and fill.fore_color.rgb:
                    print(f"Background Color: #{fill.fore_color.rgb}")
                    slide_info["background_color"] = str(fill.fore_color.rgb)
        
        # シェイプを解析
        print(f"Shapes: {len(slide.shapes)}")
        for shape in slide.shapes:
            shape_info = analyze_shape(shape)
            slide_info["shapes"].append(shape_info)
        
        all_slides_info.append(slide_info)
        print()
    
    # JSON形式でレポートを保存
    report_path = Path(pptx_path).with_suffix('.verification.json')
    with open(report_path, 'w', encoding='utf-8') as f:
        json.dump(all_slides_info, f, indent=2, ensure_ascii=False)
    print(f"Verification report saved to: {report_path}")
    
    return all_slides_info

def main():
    if len(sys.argv) < 2:
        print("Usage: python verify-pptx.py <pptx-file>")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    if not Path(pptx_path).exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)
    
    verify_pptx(pptx_path)

if __name__ == "__main__":
    main()